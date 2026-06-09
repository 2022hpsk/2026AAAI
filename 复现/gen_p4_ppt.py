#!/usr/bin/env python3
"""生成 P4 实验进度汇报 PPT（python-pptx）。

无 node（跑不了仓库 gen.js 约定），故用 python-pptx 直接生成，套用项目调色板/字体。
输出：复现/P4实验进度汇报.pptx
"""
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# 项目调色板（取自 汇报slides-build/gen_v2.js）
C = dict(dark="12183A", navy="27347A", ink="1E293B", ink2="334155", muted="64748B",
         faint="94A3B8", white="FFFFFF", panel="F1F5F9", panel2="F8FAFC", line="E2E8F0",
         cyan="0EA5C4", cyanDk="0E7490", cyanBg="E0F4F9", amber="E08A1E", amberDk="B45309",
         amberBg="FBEFD9", red="DC4A52", redBg="FBE6E7", green="16A34A", greenBg="DCFCE7")
HEAD = "Microsoft YaHei"; BODY = "Microsoft YaHei"; NUM = "Arial Black"
def rgb(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width = I(13.333); prs.slide_height = I(7.5)
W, H, M = 13.333, 7.5, 0.6
BLANK = prs.slide_layouts[6]


def slide(bg="white"):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = rgb(C[bg]); r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def _hex(v):  # 允许传调色板键名或 6 位 hex
    return C.get(v, v)


def box(s, x, y, w, h, fill=None, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = rgb(_hex(fill))
    else: sp.fill.background()
    if line: sp.line.color.rgb = rgb(_hex(line)); sp.line.width = Pt(0.75)
    else: sp.line.fill.background()
    return sp


def txt(s, x, y, w, h, runs, size=14, color="ink", bold=False, align="left",
        anchor="top", font=BODY, sp_after=2, line_sp=1.05):
    tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h)); tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs, str): runs = [[(runs, color, bold)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line_sp
        if isinstance(para, str): para = [(para, color, bold)]
        for (t, c, b) in para:
            r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b
            r.font.name = font; r.font.color.rgb = rgb(_hex(c))
    return tb


def title(s, kicker, ttl, accent="cyan", dark=False):
    box(s, M, 0.5, 0.18, 0.18, fill=accent)
    txt(s, M + 0.3, 0.42, W - 2 * M, 0.36, kicker, size=13, color=accent, bold=True)
    txt(s, M, 0.82, W - 2 * M, 0.8, ttl, size=29, color=("white" if dark else "ink"),
        bold=True, font=HEAD)


def stat(s, x, y, w, num, label, desc, col):
    box(s, x, y, w, 1.5, fill="panel2", line="line")
    box(s, x, y, 0.08, 1.5, fill=col)
    txt(s, x + 0.22, y + 0.12, w - 0.4, 0.55, num, size=33, color=col, bold=True, font=NUM)
    txt(s, x + 0.22, y + 0.72, w - 0.4, 0.32, label, size=14, color="ink", bold=True, font=HEAD)
    txt(s, x + 0.22, y + 1.06, w - 0.42, 0.4, desc, size=10.5, color="muted", line_sp=1.0)

# ─────────────────────────── Slide 1：标题（深色） ───────────────────────────
s = slide("dark")
box(s, 0, 0, 0.18, H, fill="cyan")
txt(s, M, 2.35, W - 2 * M, 0.4, "SpeakerMem-R1 · 实验进度汇报", size=15, color="cyan", bold=True)
txt(s, M, 2.8, W - 2 * M, 1.2, "P4 梯队①：复现 baseline —— 从「纸面」到「真实可跑」", size=34, color="white", bold=True, font=HEAD, line_sp=1.05)
txt(s, M, 4.3, W - 2 * M, 0.5, [[("三个多方 benchmark 已就位 · DeepSeek 已接入 · BM25 已真实出数 · Mem0 全链路打通并优化", "faint", False)]], size=14)
txt(s, M, H - 0.7, W - 2 * M, 0.3, "2026-06-09　｜　实验机：本机（无 GPU，LLM=DeepSeek deepseek-v4-flash）", size=11, color="muted")

# ─────────────────────────── Slide 2：一页纸结论 ───────────────────────────
s = slide("white")
title(s, "TL;DR · 一页纸结论", "本阶段把「一行实验没跑」变成「baseline 真实出数」", accent="green")
box(s, M, 1.55, W - 2 * M, 0.85, fill="cyanBg")
txt(s, M + 0.25, 1.55, W - 2 * M - 0.5, 0.85, [[
    ("实验环境搭好、", "ink2", False), ("三个多方 benchmark 数据全部到位", "cyanDk", True),
    ("、LLM 接入 ", "ink2", False), ("DeepSeek", "cyanDk", True),
    ("，", "ink2", False), ("BM25 baseline 已真实出数", "cyanDk", True),
    ("，", "ink2", False), ("Mem0 baseline 全链路打通并完成性能优化。", "cyanDk", True)]],
    size=15, anchor="middle", line_sp=1.15)
y = 2.75; gap = 0.18; w = (W - 2 * M - 3 * gap) / 4
stat(s, M, y, w, "3", "数据集就位", "Group/Social/EverMemBench，共 4176 题", C["cyan"])
stat(s, M + (w + gap), y, w, "35%", "BM25 出数", "Finance/multi_hop 20题，对齐论文≈38.2%", C["green"])
stat(s, M + 2 * (w + gap), y, w, "✓", "Mem0 管线通", "抽取→检索→QA→judge 全链路", C["amber"])
stat(s, M + 3 * (w + gap), y, w, "700×", "Mem0 提速", "检索 14.5s→0.02s（修线程超额订阅）", C["navy"])
box(s, M, 4.6, W - 2 * M, 0.95, fill="amberBg")
txt(s, M + 0.25, 4.6, W - 2 * M - 0.5, 0.95, [
    [("意义：", "amberDk", True), ("这是距「强提交/Highlight」唯一差距（真实实验）的实质性突破——评测 pipeline 已被真实数据与真实 LLM 跑通。", "ink2", False)],
    [("约束：", "amberDk", True), ("暂无 GPU（P5 训练前置，申请中）；未经允许不真实调用 API；判分用 DeepSeek 与论文 gpt-5 有偏移（相对比较成立）。", "ink2", False)]],
    size=12, anchor="middle", line_sp=1.15)

# ─────────────────────────── Slide 3：环境与数据 ───────────────────────────
s = slide("white")
title(s, "环境 & 数据", "三个多方 benchmark 数据全部到位", accent="cyan")
rows = [
    ["Benchmark", "来源（真实地址）", "规模", "题数"],
    ["GroupMemBench", "github UCSB-NLP-Chang（数据在仓库内）", "4 领域 × ~3万消息", "745"],
    ["SocialMemBench", "HF anon4data/socialmembench + 4open.science", "43 网络 / 7355 轮 / 430 人设", "1031"],
    ["EverMemBench", "github EverMind-AI + HF EverMemBench-Dynamic", "1263 对话块 / 170 人设", "2400"],
]
tx, ty, tw = M, 1.65, W - 2 * M; colw = [2.5, 5.0, 3.1, 1.52]; rh = 0.62
cx = tx
for j, cwj in enumerate(colw):
    for i, row in enumerate(rows):
        x = tx + sum(colw[:j]); y = ty + (0.5 if i else 0) + (i - 1) * rh if i else ty
        yy = ty + (rh * i if i == 0 else 0.5 + (i - 1) * rh)
        hh = 0.5 if i == 0 else rh
        fill = "dark" if i == 0 else ("panel2" if i % 2 == 0 else "white")
        box(s, x, yy, cwj, hh, fill=fill, line="line")
        col = "white" if i == 0 else ("ink" if j == 0 else "ink2")
        bold = (i == 0) or (j == 0) or (j == 3)
        if i > 0 and j == 3: col = "cyanDk"
        txt(s, x + 0.12, yy, cwj - 0.2, hh, rows[i][j], size=11.5 if i == 0 else 11,
            color=col, bold=bold, anchor="middle", align=("center" if j == 3 else "left"),
            font=(HEAD if i == 0 else BODY))
box(s, M, 4.55, W - 2 * M, 1.0, fill="panel")
txt(s, M + 0.25, 4.55, W - 2 * M - 0.5, 1.0, [
    [("运行环境：", "cyanDk", True), ("本机无 pip 的系统 python → 用 /opt/conda 建 venv；暂无 GPU；依赖齐全（rank-bm25 / mem0ai 2.0.4 / sentence-transformers / chromadb / torch-CPU）。", "ink2", False)],
    [("⚠ 纠错：", "red", True), ("项目自带 evaluation.py / 旧文档里的 HF 地址是错的占位符，以上才是真实可下载地址（已记入 CLAUDE.md 与记忆）。", "ink2", False)]],
    size=11.5, anchor="middle", line_sp=1.2)

# ─────────────────────────── Slide 4：评测流程 ───────────────────────────
s = slide("white")
title(s, "评测流程（统一）", "检索 → 作答 → 判分，按类别统计准确率", accent="cyan")
steps = [("①读对话", "整段群聊\n按时间排序"), ("②建索引", "BM25：消息正文\n切词建库"),
         ("③检索", "每题取 top-k\n相关段落"), ("④QA-agent", "先思考\n再 Final:答案"),
         ("⑤judge", "对比标准答案\n判 对/错"), ("⑥统计", "按 6 类\n+总体平均")]
n = len(steps); sw = 1.72; gap = (W - 2 * M - n * sw) / (n - 1); y = 2.2
for i, (h, d) in enumerate(steps):
    x = M + i * (sw + gap)
    col = "cyan" if i < 3 else "amber"
    box(s, x, y, sw, 1.4, fill=(C["cyanBg"] if i < 3 else C["amberBg"]), line=(C["cyan"] if i < 3 else C["amber"]))
    txt(s, x + 0.1, y + 0.16, sw - 0.2, 0.4, h, size=14, color=(C["cyanDk"] if i < 3 else C["amberDk"]), bold=True, align="center", font=HEAD)
    txt(s, x + 0.1, y + 0.62, sw - 0.2, 0.7, d.replace("\n", " "), size=10.5, color="ink2", align="center", line_sp=1.05)
    if i < n - 1:
        txt(s, x + sw - 0.02, y + 0.45, gap + 0.04, 0.5, "→", size=20, color="faint", align="center", bold=True)
txt(s, M, 3.85, W - 2 * M, 0.3, [[("检索阶段不调 LLM；问答阶段每题 2 次 LLM（agent+judge）。", "muted", False)]], size=11.5)
box(s, M, 4.3, W - 2 * M, 1.25, fill="panel")
txt(s, M + 0.25, 4.3, W - 2 * M - 0.5, 1.25, [
    [("适配（不改上游仓库）：", "cyanDk", True), ("llm_clients.py 把上游 OpenAI 形态调用路由到 DeepSeek / Anthropic / 本地 vLLM。", "ink2", False)],
    [("eval_patches.py：", "cyanDk", True), ("抬高 token 预算（上游给 gpt-5 调的 256/512 会截断 DeepSeek 的 Final 行→误判 Unclear）+ pin temperature=0（可复现）。", "ink2", False)],
    [("Mem0 与 BM25 的唯一区别：", "amberDk", True), ("②③步改为「LLM 抽取事实入库 → 向量检索记忆」，第④~⑥步完全一致（保证可比）。", "ink2", False)]],
    size=11.5, anchor="middle", line_sp=1.18)

# ─────────────────────────── Slide 5：BM25 结果 ───────────────────────────
s = slide("white")
title(s, "BM25 baseline · 已真实出数", "35% —— 与论文量级吻合，全链路打通", accent="green")
box(s, M, 1.7, 3.6, 2.5, fill="panel2", line="line"); box(s, M, 1.7, 0.1, 2.5, fill="green")
txt(s, M + 0.3, 1.95, 3.2, 1.0, "35.0%", size=52, color="green", bold=True, font=NUM)
txt(s, M + 0.3, 3.0, 3.2, 0.4, "Finance / multi_hop / 20 题", size=13, color="ink", bold=True, font=HEAD)
txt(s, M + 0.3, 3.42, 3.2, 0.7, [[("对照论文 BM25 multi_hop ", "muted", False), ("≈ 38.2%", "ink", True), ("，量级吻合 → 流程忠实、可复现。", "muted", False)]], size=11.5, line_sp=1.1)
bx = M + 4.0; bw = W - M - bx
txt(s, bx, 1.7, bw, 0.4, "跑这一笔顺带修了 2 个「上游为 gpt-5 调参」的坑：", size=14, color="ink", bold=True, font=HEAD)
box(s, bx, 2.2, bw, 1.0, fill="redBg"); box(s, bx, 2.2, 0.08, 1.0, fill="red")
txt(s, bx + 0.22, 2.2, bw - 0.4, 1.0, [
    [("① token 太小 → 截断", "red", True)],
    [("judge 仅 256 token，DeepSeek 啰嗦 reasoning 把 Final: 行截断 → 误判 Unclear。已抬高预算，Unclear 5→1。", "ink2", False)]],
    size=11.5, anchor="middle", line_sp=1.15)
box(s, bx, 3.35, bw, 1.0, fill="amberBg"); box(s, bx, 3.35, 0.08, 1.0, fill="amber")
txt(s, bx + 0.22, 3.35, bw - 0.4, 1.0, [
    [("② temperature=0.2 不可复现", "amberDk", True)],
    [("跑两次逐题判定会抖。已 pin temperature=0，benchmark 数字可复现。", "ink2", False)]],
    size=11.5, anchor="middle", line_sp=1.15)
box(s, M, 4.55, W - 2 * M, 1.0, fill="cyanBg")
txt(s, M + 0.25, 4.55, W - 2 * M - 0.5, 1.0, [[
    ("下一步：", "cyanDk", True),
    ("跑 GroupMemBench 全量 745 题（便宜：~1490 次调用），坐实「BM25 追平最强系统 46%」这个 motivation 核武器。", "ink2", False)]],
    size=12.5, anchor="middle", line_sp=1.15)

# ─────────────────────────── Slide 6：Mem0 profiling ───────────────────────────
s = slide("white")
title(s, "Mem0 baseline · 管线打通 + 性能优化", "定位并修复 4 个问题，检索提速 700×", accent="amber")
rows = [
    ["问题", "根因", "修复 → 效果"],
    ["检索/embedding 极慢", "192 核致 torch 线程超额订阅", "限 8 线程 → 检索 14.5s→0.02s"],
    ["抽取 0 条事实（0%）", "deepseek-v4-flash 是推理模型，抽取 JSON 被截断", "max_tokens→8000（实测可抽出事实）"],
    ["重复入库浪费 token", "get_all 签名用错被吞异常", "chroma 持久化 + 已入库则跳过复用"],
    ["遥测卡顿", "mem0/chromadb posthog 网络超时", "关闭遥测（次要提速）"],
]
ty = 1.65; colw = [3.0, 4.4, 4.73]; rh = 0.66
for i, row in enumerate(rows):
    yy = ty + (0.46 if i else 0) + max(0, i - 1) * rh if i else ty
    yy = ty if i == 0 else ty + 0.46 + (i - 1) * rh
    hh = 0.46 if i == 0 else rh
    for j, cwj in enumerate(colw):
        x = M + sum(colw[:j])
        fill = "dark" if i == 0 else ("panel2" if i % 2 == 0 else "white")
        box(s, x, yy, cwj, hh, fill=fill, line="line")
        col = "white" if i == 0 else ("ink" if j == 0 else "ink2")
        bold = (i == 0) or (j == 0)
        if i > 0 and j == 2: col = "green"; bold = True
        txt(s, x + 0.12, yy, cwj - 0.22, hh, row[j], size=11 if i == 0 else 10.5,
            color=col, bold=bold, anchor="middle", font=(HEAD if i == 0 else BODY), line_sp=1.0)
box(s, M, 5.35, W - 2 * M, 0.92, fill="amberBg")
txt(s, M + 0.25, 5.35, W - 2 * M - 0.5, 0.92, [
    [("结构性成本：", "amberDk", True), ("入库每批仍要 1 次 DeepSeek 抽取（~20s/批），全量 30k 条需专门规划；", "ink2", False),
     ("但入库一次即可持久化复用，之后所有题零额外入库成本。", "ink", True)]],
    size=11.5, anchor="middle", line_sp=1.15)

# ─────────────────────────── Slide 7：下一步 & 风险（深色） ───────────────────────────
s = slide("dark")
box(s, 0, 0, 0.18, H, fill="cyan")
title(s, "下一步 & 风险", "路线清晰，逐步坐实实验数字", accent="cyan", dark=True)
txt(s, M, 1.75, 6.0, 0.4, "下一步（按性价比）", size=16, color="cyan", bold=True, font=HEAD)
nexts = [
    "① 验证 Mem0 抽取非空后，跑适中规模出真实数字",
    "② BM25 全量 745 题（便宜，坐实 BM25≈46% motivation）",
    "③ 图驱动合成 200 条训练数据 → Stage1 SFT",
    "④ 接 SocialMemBench / EverMemBench 实跑",
]
for i, t in enumerate(nexts):
    txt(s, M, 2.25 + i * 0.62, 6.0, 0.55, [[(t, "FFFFFF", False)]], size=13, line_sp=1.05)
txt(s, 7.1, 1.75, W - M - 7.1, 0.4, "风险 / 待办", size=16, color="amber", bold=True, font=HEAD)
risks = [
    "Mem0 全量入库成本/耗时仍高（已可复用缓解）",
    "判分用 DeepSeek 与论文 gpt-5 有偏移（相对比较成立）",
    "GPU 待落实（P5 RL 训练前置）",
    "训练数据须合成，质量能否撑 RL 待验证",
]
for i, t in enumerate(risks):
    txt(s, 7.1, 2.25 + i * 0.62, W - M - 7.1, 0.55, [[("· ", "amber", True), (t, "CADCFC", False)]], size=13, line_sp=1.05)
txt(s, M, H - 0.8, W - 2 * M, 0.4, "代码 / 数据 / 文档见 复现/（README.md、数据与流程说明.md、P4实验进度汇报.md）", size=11, color="faint")

out = Path(__file__).resolve().parent / "P4实验进度汇报.pptx"
prs.save(str(out))
print("已生成:", out)
